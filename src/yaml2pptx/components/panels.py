"""Two-panel and comparison slide renderers."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from yaml2pptx.components.base import (
    add_footer,
    add_header,
    add_multiline_textbox,
    add_rect,
    add_slide_title,
    add_textbox,
)
from yaml2pptx.themes import Theme


def render_two_panels(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    left_panel: dict | None = None,
    right_panel: dict | None = None,
    **kwargs,
) -> None:
    """Render a two-panel slide (A/B comparison with different backgrounds)."""
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle, title_size=30, title_top=0.80)

    panel_top = 2.40
    panel_width = 6.05
    panel_height = 4.30
    gap = 0.30
    inner_pad = 0.25

    for i, panel in enumerate([left_panel, right_panel]):
        if not panel:
            continue

        left = s.margin_left + i * (panel_width + gap)
        is_dark = panel.get("dark", i == 0)
        bg_color = c.dark_navy if is_dark else c.card_bg

        # Panel background
        add_rect(slide, left, panel_top, panel_width, panel_height, fill=bg_color)

        inner_left = left + inner_pad
        inner_width = panel_width - inner_pad * 2

        text_color = c.white if is_dark else c.text_dark
        muted_color = c.light_blue if is_dark else c.text_muted
        label_color = c.accent if is_dark else c.primary

        # Big letter (A/B)
        if panel.get("letter"):
            letter_color = c.light_blue if is_dark else c.accent
            add_textbox(slide, inner_left, panel_top + 0.15, 1.2, 1.40,
                        text=panel["letter"], font_size=68, bold=True, color=letter_color)

        # Label
        if panel.get("label"):
            add_textbox(slide, inner_left, panel_top + 1.10, inner_width, 0.35,
                        text=panel["label"].upper(), font_size=11, bold=True, color=label_color)

        # Panel title
        if panel.get("title"):
            add_textbox(slide, inner_left, panel_top + 1.45, inner_width, 0.50,
                        text=panel["title"], font_size=20, bold=True, color=text_color)

        # Example text (italic)
        if panel.get("example"):
            add_textbox(slide, inner_left, panel_top + 1.95, inner_width, 0.50,
                        text=panel["example"], font_size=12, italic=True, color=muted_color)

        # Points
        if panel.get("points"):
            points_top = panel_top + 2.60
            points_height = panel_top + panel_height - points_top - 0.10
            font = 12 if len(panel["points"]) > 5 else 13
            add_multiline_textbox(slide, inner_left, points_top, inner_width, points_height,
                                  lines=panel["points"], font_size=font, color=text_color)

    add_footer(slide, theme)


def render_comparison(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    left_panel: dict | None = None,
    right_panel: dict | None = None,
    footer_text: str = "",
    **kwargs,
) -> None:
    """Render a comparison slide with key-value rows in two panels."""
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle, title_size=32, title_top=0.80)

    panel_top = 2.35
    panel_width = 6.05
    panel_height = 4.20
    gap = 0.30

    for i, panel in enumerate([left_panel, right_panel]):
        if not panel:
            continue

        left = s.margin_left + i * (panel_width + gap)

        # Panel background
        add_rect(slide, left, panel_top, panel_width, panel_height, fill=c.card_bg)

        # Header bar
        header_color = c.primary if i == 0 else c.accent
        add_rect(slide, left, panel_top, panel_width, 0.55, fill=header_color)

        # Header text
        if panel.get("header"):
            add_textbox(slide, left + 0.35, panel_top + 0.10, 5.0, 0.35,
                        text=panel["header"], font_size=15, bold=True, color=c.white)

        # Panel title
        if panel.get("title"):
            add_textbox(slide, left + 0.25, panel_top + 0.75, 5.55, 0.40,
                        text=panel["title"], font_size=18, bold=True, color=c.text_dark)

        # Key-value rows
        rows = panel.get("rows", [])
        label_color = header_color
        row_top = panel_top + 1.25
        available_for_rows = panel_height - 1.25 - 0.15
        row_spacing = min(0.65, max(0.40, available_for_rows / max(len(rows), 1)))
        row_font = 12 if len(rows) > 5 else 13

        for j, row in enumerate(rows):
            y = row_top + j * row_spacing
            # Label
            add_textbox(slide, left + 0.25, y, 1.40, 0.30,
                        text=row.get("label", ""), font_size=10, bold=True, color=label_color)
            # Value
            add_textbox(slide, left + 1.70, y, 4.05, row_spacing - 0.05,
                        text=row.get("value", ""), font_size=row_font, color=c.text_dark)

    # Footer text (centered, italic)
    if footer_text:
        add_textbox(slide, s.margin_left, 6.80, s.content_width, 0.30,
                    text=footer_text, font_size=12, italic=True, color=c.text_muted,
                    alignment=PP_ALIGN.CENTER)

    add_footer(slide, theme)


def render_section_divider(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    number: str = "",
    title: str = "",
    subtitle: str = "",
    **kwargs,
) -> None:
    """Render a section divider slide with large number and title."""
    c = theme.colors
    s = theme.sizes

    # Full dark background
    add_rect(slide, 0, 0, 13.34, 7.50, fill=c.dark_navy)

    # Number (large, light blue)
    if number:
        add_textbox(slide, 1.0, 1.2, 3.5, 2.00,
                    text=str(number), font_size=96, bold=True, color=c.light_blue)

    # Title
    add_textbox(slide, 1.0, 3.4, 11.0, 1.0,
                text=title, font_size=40, bold=True, color=c.white)

    # Subtitle
    if subtitle:
        add_textbox(slide, 1.0, 4.5, 11.0, 0.80,
                    text=subtitle, font_size=18, color=c.light_blue)

    # Accent line
    add_rect(slide, 1.0, 5.40, 2.0, 0.04, fill=c.accent)
