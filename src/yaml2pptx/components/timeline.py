"""Timeline and process flow slide renderers."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from yaml2pptx.components.base import (
    add_footer,
    add_header,
    add_multiline_textbox,
    add_rect,
    add_slide_title,
    add_textbox,
)
from yaml2pptx.components.icons import add_icon_shape
from yaml2pptx.themes import Theme


def render_timeline(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    phases: list[dict] | None = None,
    **kwargs,
) -> None:
    """Render a horizontal timeline with phases/milestones.

    Each phase: {label, title, description, color?, active?}
    """
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    if not phases:
        add_footer(slide, theme)
        return

    count = len(phases)
    timeline_top = 3.00
    line_y = timeline_top + 0.30
    total_width = 11.50
    phase_width = total_width / count
    start_left = 0.90

    # Horizontal timeline line
    add_rect(slide, start_left, line_y, total_width, 0.06, fill=c.text_light)

    for i, phase in enumerate(phases):
        left = start_left + i * phase_width
        is_active = phase.get("active", False)
        dot_color = c.accent if is_active else c.primary

        # Timeline dot
        dot_size = 0.30
        add_rect(slide, left + phase_width / 2 - dot_size / 2,
                 line_y - dot_size / 2 + 0.03, dot_size, dot_size, fill=dot_color)

        # Phase label (above line)
        if phase.get("label"):
            add_textbox(slide, left, timeline_top - 0.70, phase_width, 0.30,
                        text=phase["label"], font_size=10, bold=True,
                        color=c.accent if is_active else c.text_muted,
                        alignment=PP_ALIGN.CENTER)

        # Phase title (below line)
        if phase.get("title"):
            title_color = c.text_dark if is_active else c.text_dark
            add_textbox(slide, left, line_y + 0.50, phase_width, 0.45,
                        text=phase["title"], font_size=14, bold=True,
                        color=title_color, alignment=PP_ALIGN.CENTER)

        # Description
        if phase.get("description"):
            add_textbox(slide, left, line_y + 1.00, phase_width, 1.00,
                        text=phase["description"], font_size=11,
                        color=c.text_muted, alignment=PP_ALIGN.CENTER)

        # Items list
        if phase.get("items"):
            add_multiline_textbox(
                slide, left + 0.15, line_y + 1.00, phase_width - 0.30, 2.00,
                lines=phase["items"], font_size=10, color=c.text_dark, bullet="•  ",
            )

    add_footer(slide, theme)


def render_process(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    steps: list[dict] | None = None,
    **kwargs,
) -> None:
    """Render a numbered process flow with connected steps.

    Each step: {number?, title, description, icon?}
    """
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    if not steps:
        add_footer(slide, theme)
        return

    count = len(steps)
    total_width = 12.13
    gap = 0.20
    step_width = (total_width - gap * (count - 1)) / count
    step_top = 2.40
    step_height = 4.10
    inner_pad = 0.25
    # Scale font sizes for narrow steps
    title_font = 14 if count > 3 else 16
    desc_font = 9 if count > 3 else 11
    item_font = 8 if count > 3 else 9

    for i, step in enumerate(steps):
        left = s.margin_left + i * (step_width + gap)

        # Step background
        add_rect(slide, left, step_top, step_width, step_height, fill=c.card_bg)

        # Top accent line
        add_rect(slide, left, step_top, step_width, 0.06, fill=c.accent)

        inner_left = left + inner_pad
        inner_width = step_width - inner_pad * 2

        # Step number circle
        number = str(step.get("number", i + 1))
        add_rect(slide, inner_left, step_top + 0.25, 0.55, 0.55, fill=c.primary)
        add_textbox(slide, inner_left, step_top + 0.30, 0.55, 0.45,
                    text=number, font_size=22, bold=True, color=c.white,
                    alignment=PP_ALIGN.CENTER)

        # Icon (optional, beside number)
        if step.get("icon"):
            add_icon_shape(slide, inner_left + 0.65, step_top + 0.25, 0.55,
                           icon_name=step["icon"], bg_color=c.accent)

        # Step title
        if step.get("title"):
            add_textbox(slide, inner_left, step_top + 1.00, inner_width, 0.45,
                        text=step["title"], font_size=title_font, bold=True, color=c.text_dark)

        # Description
        if step.get("description"):
            add_textbox(slide, inner_left, step_top + 1.50, inner_width, 1.20,
                        text=step["description"], font_size=desc_font, color=c.text_muted)

        # Sub-items — fill remaining step space
        if step.get("items"):
            items_top = 2.70
            items_height = step_height - items_top - 0.10
            add_multiline_textbox(
                slide, inner_left, step_top + items_top, inner_width, items_height,
                lines=step["items"], font_size=item_font, color=c.text_dark, bullet="•  ",
            )

        # Arrow between steps (except last)
        if i < count - 1:
            arrow_left = left + step_width + gap / 2 - 0.10
            add_textbox(slide, arrow_left, step_top + step_height / 2 - 0.20,
                        0.20, 0.40, text="→", font_size=20, bold=True,
                        color=c.accent, alignment=PP_ALIGN.CENTER)

    add_footer(slide, theme)
