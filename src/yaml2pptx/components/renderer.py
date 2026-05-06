"""Component-based slide renderer.

This renderer builds slides from free-form shapes using a theme,
instead of relying on PowerPoint placeholder layouts.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from yaml2pptx.components.agenda import render_agenda
from yaml2pptx.components.cards import (
    render_content_cards,
    render_definition_cards,
    render_icon_cards,
    render_stat_cards,
)
from yaml2pptx.components.panels import (
    render_comparison,
    render_section_divider,
    render_two_panels,
)
from yaml2pptx.components.quote import (
    render_checklist,
    render_key_metrics,
    render_quote,
)
from yaml2pptx.components.timeline import render_process, render_timeline
from yaml2pptx.components.title_page import render_title_page
from yaml2pptx.components.base import (
    add_footer,
    add_header,
    add_multiline_textbox,
    add_rect,
    add_slide_title,
    add_textbox,
)
from yaml2pptx.themes import Theme, get_theme


def render_presentation(
    slides_data: list[dict[str, Any]],
    theme: Theme | None = None,
    output_path: str | Path = "output.pptx",
    metadata: dict[str, str] | None = None,
) -> Path:
    """Render a complete presentation from slide definitions.

    Args:
        slides_data: List of slide dicts, each with a 'type' key.
        theme: Theme to use. Defaults to the default theme.
        output_path: Output file path.
        metadata: Presentation metadata (title, author, subject).
    """
    if theme is None:
        theme = get_theme()

    prs = Presentation()
    prs.slide_width = Emu(theme.sizes.slide_width)
    prs.slide_height = Emu(theme.sizes.slide_height)

    # Use the blank layout — prefer name-based lookup, fall back to index 6
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6]

    total = len(slides_data)

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_data.get("type", "content")

        # Auto-generate page number if not specified
        if "page" not in slide_data and slide_type != "title_page":
            slide_data["page"] = f"{i + 1} / {total}"

        # Inject theme
        renderer = RENDERERS.get(slide_type)
        if renderer:
            # Filter out 'type' and 'speaker_notes' keys, pass rest as kwargs
            kwargs = {k: v for k, v in slide_data.items() if k not in ("type", "speaker_notes")}
            renderer(slide, theme, **kwargs)
        else:
            # Fallback: simple content slide
            _render_simple_content(slide, theme, slide_data)

        # Speaker notes (supported on all slide types)
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

    # Metadata
    if metadata:
        if metadata.get("title"):
            prs.core_properties.title = metadata["title"]
        if metadata.get("author"):
            prs.core_properties.author = metadata["author"]
        if metadata.get("subject"):
            prs.core_properties.subject = metadata["subject"]

    out = Path(output_path)
    prs.save(str(out))
    return out


def _add_image(slide, image_path: str, left: float, top: float,
               width: float | None = None, height: float | None = None) -> None:
    """Add an image to a slide, with optional sizing."""
    path = Path(image_path)
    if not path.exists():
        raise FileNotFoundError(f"Image not found: {path}")
    img_width = Inches(width) if width else None
    img_height = Inches(height) if height else None
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), img_width, img_height)


def render_image_text(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    image: str = "",
    image_width: float = 5.5,
    image_position: str = "left",
    content: list | str = "",
    caption: str = "",
    **kwargs,
) -> None:
    """Render a slide with image and text side by side.

    Args:
        image: Path to image file.
        image_width: Width of image in inches.
        image_position: "left" or "right".
        content: Text content (list of bullets or string).
        caption: Optional image caption.
    """
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    content_top = s.content_top
    gap = 0.40
    text_width = s.content_width - image_width - gap

    if image_position == "right":
        text_left = s.margin_left
        img_left = s.margin_left + text_width + gap
    else:
        img_left = s.margin_left
        text_left = s.margin_left + image_width + gap

    # Image
    if image:
        try:
            _add_image(slide, image, img_left, content_top, width=image_width)
        except FileNotFoundError:
            logger.warning("Image not found: %s", image)
            # Render placeholder rectangle if image not found
            add_rect(slide, img_left, content_top, image_width, 3.5, fill=c.card_bg)
            add_textbox(slide, img_left, content_top + 1.5, image_width, 0.5,
                        text=f"[Bilde: {image}]", font_size=11, color=c.text_muted)

    # Caption
    if caption:
        add_textbox(slide, img_left, content_top + 3.7, image_width, 0.3,
                    text=caption, font_size=9, italic=True, color=c.text_muted)

    # Content
    if content:
        if isinstance(content, list):
            add_multiline_textbox(
                slide, text_left, content_top, text_width, 4.0,
                lines=content, font_size=s.body_size, color=c.text_dark,
                bullet="•  ",
            )
        elif isinstance(content, str):
            add_textbox(slide, text_left, content_top, text_width, 4.0,
                        text=content, font_size=s.body_size, color=c.text_dark)

    add_footer(slide, theme)


def _add_themed_table(
    slide,
    theme: Theme,
    headers: list[str] | None,
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    row_height: float = 0.40,
) -> None:
    """Add a styled table to a slide."""
    c = theme.colors
    s = theme.sizes
    has_headers = headers is not None
    if not rows and not has_headers:
        return
    num_rows = len(rows) + (1 if has_headers else 0)
    num_cols = len(headers) if has_headers else (len(rows[0]) if rows else 0)
    if num_cols == 0:
        return

    # Auto-scale row height to fit within slide
    max_height = s.footer_top - top - 0.15
    total_height = row_height * num_rows
    if total_height > max_height:
        row_height = max_height / num_rows
    height = row_height * num_rows

    # Use smaller font for dense tables
    font_size = 11 if row_height >= 0.35 else 9

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    table = table_shape.table

    row_offset = 0
    if has_headers and headers:
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = True
                    run.font.name = theme.fonts.primary
        row_offset = 1

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + row_offset, col_idx)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = theme.fonts.primary


def render_table(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    headers: list[str] | None = None,
    rows: list[list[str]] | None = None,
    **kwargs,
) -> None:
    """Render a slide with a data table."""
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    if rows or headers:
        _add_themed_table(
            slide, theme,
            headers=headers,
            rows=rows or [],
            left=s.margin_left,
            top=s.content_top,
            width=s.content_width,
        )

    add_footer(slide, theme)


def _render_simple_content(slide, theme: Theme, data: dict) -> None:
    """Fallback renderer for simple content slides."""
    c = theme.colors
    s = theme.sizes

    section = data.get("section", "")
    page = data.get("page", "")
    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    content = data.get("content", [])
    image = data.get("image", "")
    table = data.get("table")

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    # Table below content
    if table and isinstance(table, dict):
        table_top = s.content_top if not content else s.content_top + 2.5
        _add_themed_table(
            slide, theme,
            headers=table.get("headers"),
            rows=table.get("rows", []),
            left=s.margin_left,
            top=table_top,
            width=s.content_width,
        )

    # If image provided, shrink content area
    if image:
        try:
            _add_image(slide, image, 8.0, s.content_top, width=4.5)
        except FileNotFoundError:
            logger.warning("Image not found: %s", image)
            add_rect(slide, 8.0, s.content_top, 4.5, 3.5, fill=c.card_bg)
        content_width = 7.0
    else:
        content_width = s.content_width

    if content:
        if isinstance(content, list):
            add_multiline_textbox(
                slide, s.margin_left, s.content_top, content_width, 4.0,
                lines=content, font_size=s.body_size, color=c.text_dark,
                bullet="•  ",
            )
        elif isinstance(content, str):
            add_textbox(slide, s.margin_left, s.content_top, content_width, 4.0,
                        text=content, font_size=s.body_size, color=c.text_dark)

    add_footer(slide, theme)


# Map of slide type -> renderer function
RENDERERS = {
    "title_page": render_title_page,
    "agenda": render_agenda,
    "stat_cards": render_stat_cards,
    "definition_cards": render_definition_cards,
    "content_cards": render_content_cards,
    "icon_cards": render_icon_cards,
    "two_panels": render_two_panels,
    "comparison": render_comparison,
    "section_divider": render_section_divider,
    "timeline": render_timeline,
    "process": render_process,
    "quote": render_quote,
    "key_metrics": render_key_metrics,
    "checklist": render_checklist,
    "image_text": render_image_text,
    "table": render_table,
}
