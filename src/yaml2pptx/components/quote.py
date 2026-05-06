"""Quote and highlight slide renderers."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN

from yaml2pptx.components.base import (
    add_footer,
    add_header,
    add_multiline_textbox,
    add_rect,
    add_slide_title,
    add_textbox,
)
from yaml2pptx.themes import Theme, resolve_color


def render_quote(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    quote: str = "",
    attribution: str = "",
    source: str = "",
    dark: bool = False,
    **kwargs,
) -> None:
    """Render a quote slide with large text and attribution.

    Args:
        quote: The quote text.
        attribution: Who said it (name/role).
        source: Where it's from (publication, year).
        dark: Use dark background.
    """
    c = theme.colors
    s = theme.sizes

    if dark:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c.dark_navy

    add_header(slide, theme, section, page)

    # Quote mark (decorative)
    quote_mark_color = c.accent if not dark else c.light_blue
    add_textbox(slide, 1.20, 1.40, 1.5, 2.0, text="\u201C",
                font_size=96, color=quote_mark_color, bold=True)

    # Quote text — scale font for long quotes
    text_color = c.text_dark if not dark else c.white
    quote_font = 22 if len(quote) > 200 else 24 if len(quote) > 120 else 26
    add_textbox(slide, 1.50, 2.30, 10.30, 2.80, text=quote,
                font_size=quote_font, color=text_color, italic=True)

    # Accent line
    add_rect(slide, 1.50, 5.20, 1.50, 0.04, fill=c.accent)

    # Attribution
    if attribution:
        attr_color = c.text_dark if not dark else c.white
        add_textbox(slide, 1.50, 5.40, 10.0, 0.35, text=attribution,
                    font_size=14, bold=True, color=attr_color)

    # Source
    if source:
        src_color = c.text_muted if not dark else c.light_blue
        add_textbox(slide, 1.50, 5.75, 10.0, 0.30, text=source,
                    font_size=11, italic=True, color=src_color)

    if not dark:
        add_footer(slide, theme)


def render_key_metrics(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    metrics: list[dict] | None = None,
    **kwargs,
) -> None:
    """Render a dashboard-style metrics grid.

    Each metric: {value, label, trend?, trend_label?, color?}
    Renders in a 2-row grid if more than 4 metrics.
    """
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    if not metrics:
        add_footer(slide, theme)
        return

    count = len(metrics)
    cols = min(count, 4)
    rows = (count + cols - 1) // cols
    total_width = 12.13
    gap = 0.20
    metric_width = (total_width - gap * (cols - 1)) / cols
    metric_height = 1.80 if rows > 1 else 3.20
    start_top = 2.20
    inner_pad = 0.25

    for idx, metric in enumerate(metrics):
        row = idx // cols
        col = idx % cols
        left = s.margin_left + col * (metric_width + gap)
        top = start_top + row * (metric_height + gap)

        # Metric background
        add_rect(slide, left, top, metric_width, metric_height, fill=c.card_bg)

        # Left accent bar
        color_name = metric.get("color")
        accent_color = resolve_color(color_name, c.primary, theme)
        add_rect(slide, left, top, 0.06, metric_height, fill=accent_color)

        inner_left = left + inner_pad
        inner_width = metric_width - inner_pad * 2

        # Value (big number)
        if metric.get("value"):
            add_textbox(slide, inner_left, top + 0.20, inner_width, 0.80,
                        text=str(metric["value"]), font_size=36, bold=True,
                        color=accent_color)

        # Label
        if metric.get("label"):
            add_textbox(slide, inner_left, top + 1.00, inner_width, 0.35,
                        text=metric["label"], font_size=12, bold=True, color=c.text_dark)

        # Trend indicator
        if metric.get("trend"):
            trend = metric["trend"]
            _trend_colors = {"up": c.success, "down": c.warning}
            trend_color = _trend_colors.get(trend, c.text_muted)
            trend_symbol = "▲" if trend == "up" else "▼" if trend == "down" else "●"
            trend_text = f"{trend_symbol}  {metric.get('trend_label', '')}"
            add_textbox(slide, inner_left, top + 1.35, inner_width, 0.30,
                        text=trend_text, font_size=10, color=trend_color)

    add_footer(slide, theme)


def render_checklist(
    slide,
    theme: Theme,
    *,
    section: str = "",
    page: str = "",
    title: str = "",
    subtitle: str = "",
    items: list[dict] | None = None,
    columns: int = 1,
    **kwargs,
) -> None:
    """Render a checklist with status indicators.

    Each item: {text, status: "done"|"in_progress"|"pending"|"blocked", note?}
    """
    c = theme.colors
    s = theme.sizes

    add_header(slide, theme, section, page)
    add_slide_title(slide, theme, title, subtitle)

    if not items:
        add_footer(slide, theme)
        return

    STATUS_MAP = {
        "done": ("✓", c.success),
        "in_progress": ("●", c.accent),
        "pending": ("○", c.text_muted),
        "blocked": ("✗", c.warning),
    }

    col_count = max(1, min(columns, 3))
    col_width = (12.13 - 0.30 * (col_count - 1)) / col_count
    start_top = 2.20
    items_per_col = (len(items) + col_count - 1) // col_count

    # Adaptive row height based on available space and whether notes exist
    has_notes = any(isinstance(it, dict) and it.get("note") for it in items)
    available = 4.50  # space before footer
    min_row = 0.55 if has_notes else 0.42
    row_height = min(0.80, max(min_row, available / items_per_col))

    for idx, item in enumerate(items):
        col = idx // items_per_col if col_count > 1 else 0
        row = idx % items_per_col if col_count > 1 else idx

        left = s.margin_left + col * (col_width + 0.30)
        top = start_top + row * row_height

        status = item.get("status", "pending")
        symbol, color = STATUS_MAP.get(status, STATUS_MAP["pending"])

        # Status indicator
        add_textbox(slide, left, top, 0.40, 0.35, text=symbol,
                    font_size=16, bold=True, color=color)

        # Item text — use font that fits row height
        text = item.get("text", "") if isinstance(item, dict) else str(item)
        text_color = c.text_dark if status != "blocked" else c.text_muted
        text_height = row_height * 0.55 if has_notes else row_height - 0.05
        add_textbox(slide, left + 0.45, top, col_width - 0.80, text_height,
                    text=text, font_size=11, color=text_color)

        # Note (small, muted) — positioned below the text
        if isinstance(item, dict) and item.get("note"):
            note_top = row_height * 0.55
            note_height = row_height - note_top - 0.03
            add_textbox(slide, left + 0.45, top + note_top, col_width - 0.80, note_height,
                        text=item["note"], font_size=9, italic=True, color=c.text_muted)

    add_footer(slide, theme)


