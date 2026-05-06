"""Base helpers for building slides from shapes."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from yaml2pptx.parser.markdown_parser import RunSegment, parse_markdown
from yaml2pptx.themes import Theme

MONO_FONT = "Consolas"


def _apply_segment(run, segment: RunSegment, font_name: str) -> None:
    """Apply Markdown segment formatting to a python-pptx run."""
    if segment.bold:
        run.font.bold = True
    if segment.italic:
        run.font.italic = True
    if segment.code:
        run.font.name = MONO_FONT
    else:
        run.font.name = font_name
    if segment.strikethrough:
        run.font.strikethrough = True
    if segment.hyperlink:
        run.hyperlink.address = segment.hyperlink


def _add_runs_to_paragraph(
    p,
    text: str,
    font_size: int,
    font_name: str,
    color: RGBColor | None,
    bold: bool,
    italic: bool,
) -> None:
    """Parse Markdown in text and add formatted runs to a paragraph."""
    segments = parse_markdown(text)

    for segment in segments:
        run = p.add_run()
        run.text = segment.text
        run.font.size = Pt(font_size)
        run.font.name = font_name
        if color:
            run.font.color.rgb = color

        # Base formatting from parameters
        run.font.bold = bold
        run.font.italic = italic

        # Markdown overrides
        _apply_segment(run, segment, font_name)


def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor | None = None,
    border: bool = False,
) -> object:
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    font_size: int = 13,
    font_name: str = "Calibri",
    color: RGBColor | None = None,
    bold: bool = False,
    italic: bool = False,
    alignment: PP_ALIGN | None = None,
    word_wrap: bool = True,
) -> object:
    """Add a text box with a single paragraph. Supports inline Markdown."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment

    _add_runs_to_paragraph(p, text, font_size, font_name, color, bold, italic)
    return txBox


def add_multiline_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str | dict],
    font_size: int = 13,
    font_name: str = "Calibri",
    color: RGBColor | None = None,
    bold: bool = False,
    italic: bool = False,
    bullet: str = "",
    line_spacing: float | None = None,
    alignment: PP_ALIGN | None = None,
) -> object:
    """Add a text box with multiple paragraphs. Supports inline Markdown and indentation levels.

    Lines can be strings or dicts with {text, level} for nested bullets.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if alignment:
            p.alignment = alignment

        # Support dict items with level
        if isinstance(line, dict):
            line_text = line.get("text", "")
            level = line.get("level", 0)
            p.level = level
        else:
            line_text = str(line)

        text = f"{bullet}{line_text}" if bullet else line_text
        _add_runs_to_paragraph(p, text, font_size, font_name, color, bold, italic)
    return txBox


def add_line(
    slide,
    left: float,
    top: float,
    width: float,
    color: RGBColor | None = None,
    weight: float = 0.75,
) -> object:
    """Add a horizontal line (connector)."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(left), Inches(top),
        Inches(left + width), Inches(top),
    )
    if color:
        connector.line.color.rgb = color
    connector.line.width = Pt(weight)
    return connector


def add_header(
    slide,
    theme: Theme,
    section: str,
    page: str,
) -> None:
    """Add the standard slide header: teal dot, section label, page number."""
    s = theme.sizes
    c = theme.colors

    # Teal dot
    add_rect(slide, s.header_dot_left, s.header_dot_top,
             s.header_dot_size, s.header_dot_size, fill=c.accent)

    # Section label
    add_textbox(slide, s.section_label_left, s.section_label_top,
                8.0, 0.30, text=section.upper(),
                font_size=s.section_label_size, bold=True, color=c.text_muted)

    # Page number
    add_textbox(slide, s.page_num_left, s.page_num_top,
                0.70, 0.30, text=page,
                font_size=s.page_num_size, color=c.text_muted,
                alignment=PP_ALIGN.RIGHT)


def add_footer(slide, theme: Theme) -> None:
    """Add the standard slide footer."""
    s = theme.sizes
    c = theme.colors
    footer = theme.footer_text
    if not footer and theme.organization:
        parts = [theme.organization]
        if theme.document_title:
            parts.append(theme.document_title)
        if theme.classification:
            parts.append(theme.classification)
        footer = "   •   ".join(parts)
    if footer:
        add_textbox(slide, s.margin_left, s.footer_top,
                    s.content_width, s.footer_height, text=footer,
                    font_size=s.footer_size, color=c.text_muted)


def add_slide_title(
    slide,
    theme: Theme,
    title: str,
    subtitle: str = "",
    title_size: int | None = None,
    title_top: float | None = None,
) -> None:
    """Add title and subtitle in the standard position."""
    s = theme.sizes
    c = theme.colors
    t_size = title_size or s.title_size
    t_top = title_top or s.title_top

    add_textbox(slide, s.margin_left, t_top,
                12.0, s.title_height, text=title,
                font_size=t_size, bold=True, color=c.text_dark)
    if subtitle:
        add_textbox(slide, s.margin_left, s.subtitle_top,
                    12.0, s.subtitle_height, text=subtitle,
                    font_size=s.subtitle_size, color=c.text_muted)
