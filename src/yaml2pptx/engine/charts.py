"""Chart generation for python-pptx slides."""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from yaml2pptx.models.elements import ChartData

CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


def add_chart_to_slide(
    slide,  # type: ignore[no-untyped-def]
    chart_data_model: ChartData,
    left: float = 1.0,
    top: float = 2.0,
    width: float = 8.0,
    height: float = 4.5,
) -> None:
    """Add a chart to a slide."""
    chart_data = CategoryChartData()
    chart_data.categories = chart_data_model.categories

    for series in chart_data_model.series:
        chart_data.add_series(series.name, series.values)

    chart_type = CHART_TYPE_MAP[chart_data_model.type]

    slide.shapes.add_chart(
        chart_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    )
