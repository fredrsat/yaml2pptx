"""Text frame and paragraph building utilities for python-pptx."""

from __future__ import annotations

from pptx.util import Pt

from yaml2pptx.models.elements import TextElement
from yaml2pptx.parser.markdown_parser import RunSegment, parse_markdown


MONO_FONT = "Consolas"


def apply_run_formatting(run, segment: RunSegment) -> None:  # type: ignore[no-untyped-def]
    """Apply formatting from a RunSegment to a python-pptx Run."""
    if segment.bold:
        run.font.bold = True
    if segment.italic:
        run.font.italic = True
    if segment.code:
        run.font.name = MONO_FONT
    if segment.strikethrough:
        run.font.strikethrough = True
    if segment.hyperlink:
        run.hyperlink.address = segment.hyperlink


def populate_text_frame(
    text_frame,  # type: ignore[no-untyped-def]
    elements: list[TextElement],
    font_size: int | None = None,
    font_name: str | None = None,
) -> None:
    """Populate a text frame with TextElements, applying Markdown formatting."""
    text_frame.clear()

    for i, element in enumerate(elements):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.level = element.level

        segments = parse_markdown(element.text)

        for j, segment in enumerate(segments):
            if j == 0:
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            else:
                run = paragraph.add_run()

            run.text = segment.text
            apply_run_formatting(run, segment)

            # Apply element-level overrides
            if element.bold is not None:
                run.font.bold = element.bold
            if element.italic is not None:
                run.font.italic = element.italic

            # Apply defaults
            size = element.font_size or font_size
            if size:
                run.font.size = Pt(size)
            if font_name and not segment.code:
                run.font.name = font_name


def set_title(slide, title_text: str) -> None:  # type: ignore[no-untyped-def]
    """Set the title of a slide, with Markdown support."""
    if not slide.shapes.title:
        return

    tf = slide.shapes.title.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]

    segments = parse_markdown(title_text)
    for i, segment in enumerate(segments):
        if i == 0:
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        else:
            run = paragraph.add_run()
        run.text = segment.text
        apply_run_formatting(run, segment)
