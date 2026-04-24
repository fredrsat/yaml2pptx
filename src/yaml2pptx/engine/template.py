from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.slide import SlideLayout


@dataclass
class PlaceholderInfo:
    idx: int
    type: str
    name: str
    width: int
    height: int


@dataclass
class LayoutInfo:
    index: int
    name: str
    placeholders: list[PlaceholderInfo] = field(default_factory=list)


class TemplateError(Exception):
    pass


class TemplateManager:
    def __init__(self, template_path: str | Path | None = None):
        if template_path:
            path = Path(template_path)
            if not path.exists():
                raise TemplateError(f"Template not found: {path}")
            self.prs = Presentation(str(path))
        else:
            self.prs = Presentation()

        self._layouts: dict[str, LayoutInfo] = {}
        self._layout_by_index: dict[int, LayoutInfo] = {}
        self._discover_layouts()

    def _discover_layouts(self) -> None:
        for i, layout in enumerate(self.prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append(
                    PlaceholderInfo(
                        idx=ph.placeholder_format.idx,
                        type=str(ph.placeholder_format.type),
                        name=ph.name,
                        width=ph.width,
                        height=ph.height,
                    )
                )
            info = LayoutInfo(index=i, name=layout.name, placeholders=placeholders)
            self._layouts[layout.name] = info
            self._layout_by_index[i] = info

    def get_layout(
        self, name: str, layout_map: dict[str, str | int] | None = None
    ) -> SlideLayout:
        """Resolve a layout name to a SlideLayout object.

        Resolution order:
        1. Explicit layout_map (if provided)
        2. Exact name match
        3. Case-insensitive name match
        4. Integer index fallback
        """
        # 1. Check layout_map
        if layout_map and name in layout_map:
            mapped = layout_map[name]
            if isinstance(mapped, int):
                return self.prs.slide_layouts[mapped]
            name = mapped

        # 2. Exact name match
        for layout in self.prs.slide_layouts:
            if layout.name == name:
                return layout

        # 3. Case-insensitive match
        name_lower = name.lower().replace("_", " ")
        for layout in self.prs.slide_layouts:
            if layout.name.lower().replace("_", " ") == name_lower:
                return layout

        # 4. Try as integer index
        try:
            idx = int(name)
            return self.prs.slide_layouts[idx]
        except (ValueError, IndexError):
            pass

        available = [f"  [{i}] {info.name}" for i, info in self._layout_by_index.items()]
        raise TemplateError(
            f"Layout '{name}' not found in template.\n"
            f"Available layouts:\n" + "\n".join(available)
        )

    def get_layouts_info(self) -> list[LayoutInfo]:
        return [self._layout_by_index[i] for i in sorted(self._layout_by_index)]

    def format_inspect(self) -> str:
        lines = ["Layouts:"]
        for info in self.get_layouts_info():
            lines.append(f"  [{info.index}] \"{info.name}\"")
            for ph in info.placeholders:
                lines.append(f"      - idx={ph.idx}, type={ph.type}, name=\"{ph.name}\"")
        return "\n".join(lines)
