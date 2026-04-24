from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from yaml2pptx.engine.generator import GenerationError, generate
from yaml2pptx.models.presentation import (
    PresentationMetadata,
    PresentationModel,
)
from yaml2pptx.models.slides import SlideDefinition


class TestGenerate:
    def test_minimal_presentation(self, tmp_path):
        model = PresentationModel(
            slides=[SlideDefinition(layout="Title Slide", title="Hello")]
        )
        out = generate(model, tmp_path / "out.pptx")
        assert out.exists()

        prs = Presentation(str(out))
        assert len(prs.slides) == 1
        assert prs.slides[0].shapes.title.text == "Hello"

    def test_content_slide(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(
                    layout="Title and Content",
                    title="Test",
                    content=["Bullet 1", "Bullet 2"],
                )
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        body = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body = shape
        assert body is not None
        assert "Bullet 1" in body.text_frame.text

    def test_speaker_notes(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(
                    layout="Title and Content",
                    title="Test",
                    speaker_notes="My notes",
                )
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "My notes" in notes

    def test_metadata(self, tmp_path):
        model = PresentationModel(
            metadata=PresentationMetadata(
                title="My Title", author="Author", subject="Subject"
            ),
            slides=[SlideDefinition(layout="Title Slide", title="Hi")],
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert prs.core_properties.title == "My Title"
        assert prs.core_properties.author == "Author"
        assert prs.core_properties.subject == "Subject"

    def test_two_column(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(
                    layout="Two Content",
                    title="Comparison",
                    left=["Left 1"],
                    right=["Right 1"],
                )
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        texts = []
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in (1, 2):
                texts.append(shape.text_frame.text)
        assert "Left 1" in texts[0]
        assert "Right 1" in texts[1]

    def test_invalid_layout(self, tmp_path):
        model = PresentationModel(
            slides=[SlideDefinition(layout="Nonexistent Layout", title="Test")]
        )
        with pytest.raises(GenerationError, match="Slide 1"):
            generate(model, tmp_path / "out.pptx")

    def test_output_from_model(self, tmp_path):
        model = PresentationModel(
            output=str(tmp_path / "model_output.pptx"),
            slides=[SlideDefinition(layout="Title Slide", title="Hi")],
        )
        out = generate(model)
        assert out == tmp_path / "model_output.pptx"
        assert out.exists()

    def test_markdown_in_content(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(
                    layout="Title and Content",
                    title="Test",
                    content=["Hello **bold** world"],
                )
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                para = shape.text_frame.paragraphs[0]
                runs = para.runs
                assert len(runs) == 3
                assert runs[1].text == "bold"
                assert runs[1].font.bold is True

    def test_table_slide(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(
                    layout="Title and Content",
                    title="Data",
                    table={
                        "headers": ["A", "B"],
                        "rows": [["1", "2"], ["3", "4"]],
                    },
                )
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        has_table = any(shape.has_table for shape in slide.shapes)
        assert has_table

    def test_chart_slide(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(
                    layout="Title and Content",
                    title="Chart",
                    chart={
                        "type": "bar",
                        "categories": ["Q1", "Q2"],
                        "series": [{"name": "Rev", "values": [100, 200]}],
                    },
                )
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        has_chart = any(shape.has_chart for shape in slide.shapes)
        assert has_chart

    def test_multiple_slides(self, tmp_path):
        model = PresentationModel(
            slides=[
                SlideDefinition(layout="Title Slide", title="Slide 1"),
                SlideDefinition(layout="Title and Content", title="Slide 2", content=["A"]),
                SlideDefinition(layout="Section Header", title="Slide 3"),
            ]
        )
        out = generate(model, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 3
