"""Edge case tests for tables and charts."""

from __future__ import annotations

import pytest
from pptx import Presentation

from yaml2pptx.engine.tables import add_table_to_slide
from yaml2pptx.engine.charts import add_chart_to_slide
from yaml2pptx.models.elements import TableData, ChartData, ChartSeries


@pytest.fixture
def slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


class TestTableEdgeCases:
    def test_empty_rows_no_headers(self, slide):
        """Empty rows with no headers should not crash."""
        table_data = TableData(rows=[])
        add_table_to_slide(slide, table_data)
        # Should return without adding anything

    def test_empty_rows_with_headers(self, slide):
        """Headers only, no data rows."""
        table_data = TableData(headers=["A", "B", "C"], rows=[])
        add_table_to_slide(slide, table_data)
        assert len(slide.shapes) == 1

    def test_normal_table(self, slide):
        table_data = TableData(
            headers=["Name", "Value"],
            rows=[["A", "1"], ["B", "2"]],
        )
        add_table_to_slide(slide, table_data)
        assert len(slide.shapes) == 1

    def test_table_no_headers(self, slide):
        table_data = TableData(rows=[["A", "1"], ["B", "2"]])
        add_table_to_slide(slide, table_data)
        assert len(slide.shapes) == 1


class TestChartEdgeCases:
    def test_basic_chart(self, slide):
        chart_data = ChartData(
            type="bar",
            categories=["Q1", "Q2"],
            series=[ChartSeries(name="S1", values=[10, 20])],
        )
        add_chart_to_slide(slide, chart_data)
        assert len(slide.shapes) == 1

    def test_invalid_chart_type(self, slide):
        """Invalid chart type should raise ValidationError from Pydantic."""
        with pytest.raises(Exception):
            ChartData(
                type="invalid",
                categories=["A"],
                series=[ChartSeries(name="S", values=[1])],
            )

    def test_pie_chart(self, slide):
        chart_data = ChartData(
            type="pie",
            categories=["A", "B", "C"],
            series=[ChartSeries(name="Data", values=[30, 50, 20])],
        )
        add_chart_to_slide(slide, chart_data)
        assert len(slide.shapes) == 1
