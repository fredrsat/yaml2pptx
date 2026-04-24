"""Tests for the component-based slide renderers (build pipeline)."""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Emu

from yaml2pptx.components.renderer import RENDERERS, render_presentation
from yaml2pptx.themes import Theme, get_theme, resolve_color


@pytest.fixture
def theme():
    return get_theme()


@pytest.fixture
def blank_slide(theme):
    """Create a blank presentation and return a single blank slide."""
    prs = Presentation()
    prs.slide_width = Emu(theme.sizes.slide_width)
    prs.slide_height = Emu(theme.sizes.slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide


# --- render_presentation ---

class TestRenderPresentation:
    def test_basic_content_slide(self, tmp_path):
        slides = [{"type": "content", "title": "Hello", "content": ["Item 1"]}]
        out = render_presentation(slides, output_path=tmp_path / "out.pptx")
        assert out.exists()

    def test_empty_slides_list(self, tmp_path):
        out = render_presentation([], output_path=tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 0

    def test_metadata_applied(self, tmp_path):
        slides = [{"type": "content", "title": "T"}]
        meta = {"title": "My Title", "author": "Author", "subject": "Sub"}
        out = render_presentation(slides, output_path=tmp_path / "out.pptx", metadata=meta)
        prs = Presentation(str(out))
        assert prs.core_properties.title == "My Title"
        assert prs.core_properties.author == "Author"

    def test_auto_page_numbers(self, tmp_path):
        slides = [
            {"type": "title_page", "title": "Title"},
            {"type": "content", "title": "S1"},
            {"type": "content", "title": "S2"},
        ]
        render_presentation(slides, output_path=tmp_path / "out.pptx")
        # title_page should not get page number
        assert "page" not in slides[0]
        assert slides[1]["page"] == "2 / 3"
        assert slides[2]["page"] == "3 / 3"

    def test_speaker_notes(self, tmp_path):
        slides = [
            {"type": "content", "title": "T", "speaker_notes": "My notes"},
        ]
        out = render_presentation(slides, output_path=tmp_path / "out.pptx")
        prs = Presentation(str(out))
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "My notes" in notes

    def test_unknown_type_uses_fallback(self, tmp_path):
        slides = [{"type": "nonexistent_type", "title": "Test", "content": ["A"]}]
        out = render_presentation(slides, output_path=tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_all_renderers_registered(self):
        expected = {
            "title_page", "agenda", "stat_cards", "definition_cards",
            "content_cards", "icon_cards", "two_panels", "comparison",
            "section_divider", "timeline", "process", "quote",
            "key_metrics", "checklist", "image_text", "table",
        }
        assert set(RENDERERS.keys()) == expected


# --- Individual renderer tests ---

class TestTitlePage:
    def test_renders_without_error(self, blank_slide, theme):
        from yaml2pptx.components.title_page import render_title_page
        render_title_page(blank_slide, theme, title="Test", subtitle="Sub",
                          author="Auth", date="2026", category="CAT")
        assert len(blank_slide.shapes) > 0

    def test_dark_background(self, blank_slide, theme):
        from yaml2pptx.components.title_page import render_title_page
        render_title_page(blank_slide, theme, title="Test")
        bg = blank_slide.background.fill
        assert bg.fore_color.rgb == theme.colors.dark_navy

    def test_extra_kwargs_ignored(self, blank_slide, theme):
        from yaml2pptx.components.title_page import render_title_page
        render_title_page(blank_slide, theme, title="T", unknown_field="ignored")


class TestAgenda:
    def test_renders_items(self, blank_slide, theme):
        from yaml2pptx.components.agenda import render_agenda
        items = [
            {"title": "Topic 1", "description": "Desc 1"},
            {"title": "Topic 2", "description": "Desc 2"},
        ]
        render_agenda(blank_slide, theme, title="Agenda", items=items)
        assert len(blank_slide.shapes) > 0

    def test_empty_items(self, blank_slide, theme):
        from yaml2pptx.components.agenda import render_agenda
        render_agenda(blank_slide, theme, title="Empty")
        # Should still render header and title

    def test_extra_kwargs_ignored(self, blank_slide, theme):
        from yaml2pptx.components.agenda import render_agenda
        render_agenda(blank_slide, theme, title="T", custom_field="x")


class TestStatCards:
    def test_renders_cards(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_stat_cards
        cards = [
            {"stat": "85%", "label": "metric", "title": "T", "description": "D"},
        ]
        render_stat_cards(blank_slide, theme, title="Stats", cards=cards)
        assert len(blank_slide.shapes) > 0

    def test_with_footnotes(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_stat_cards
        cards = [{"stat": "42"}]
        render_stat_cards(blank_slide, theme, cards=cards, footnotes=["[1] Source"])

    def test_empty_cards(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_stat_cards
        render_stat_cards(blank_slide, theme, title="Empty")


class TestDefinitionCards:
    def test_renders_with_icons(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_definition_cards
        cards = [
            {"term": "API", "icon": "gear", "subtitle": "Sub", "description": "Desc"},
        ]
        render_definition_cards(blank_slide, theme, cards=cards)
        assert len(blank_slide.shapes) > 0

    def test_renders_without_icon(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_definition_cards
        cards = [{"term": "Test"}]
        render_definition_cards(blank_slide, theme, cards=cards)

    def test_with_callout(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_definition_cards
        cards = [{"term": "X"}]
        callout = {"label": "NOTE", "columns": [{"title": "T", "description": "D"}]}
        render_definition_cards(blank_slide, theme, cards=cards, callout=callout)

    def test_custom_border_color(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_definition_cards
        cards = [{"term": "X", "border_color": "#FF0000"}]
        render_definition_cards(blank_slide, theme, cards=cards)


class TestContentCards:
    def test_renders_with_points(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_content_cards
        cards = [{"title": "T", "icon": "shield", "points": ["A", "B"]}]
        render_content_cards(blank_slide, theme, cards=cards)
        assert len(blank_slide.shapes) > 0

    def test_renders_with_description(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_content_cards
        cards = [{"title": "T", "subtitle": "S", "description": "D"}]
        render_content_cards(blank_slide, theme, cards=cards)


class TestIconCards:
    def test_renders_message_and_cards(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_icon_cards
        cards = [{"title": "T", "description": "D"}]
        render_icon_cards(blank_slide, theme, message="Key message", cards=cards)
        assert len(blank_slide.shapes) > 0

    def test_no_cards(self, blank_slide, theme):
        from yaml2pptx.components.cards import render_icon_cards
        render_icon_cards(blank_slide, theme, message="Just message")


class TestTwoPanels:
    def test_renders_both_panels(self, blank_slide, theme):
        from yaml2pptx.components.panels import render_two_panels
        render_two_panels(
            blank_slide, theme, title="T",
            left_panel={"letter": "A", "label": "LEFT", "title": "LT", "points": ["P1"]},
            right_panel={"letter": "B", "label": "RIGHT", "title": "RT", "points": ["P2"]},
        )
        assert len(blank_slide.shapes) > 0

    def test_single_panel(self, blank_slide, theme):
        from yaml2pptx.components.panels import render_two_panels
        render_two_panels(blank_slide, theme, left_panel={"title": "Only left"})


class TestComparison:
    def test_renders_with_rows(self, blank_slide, theme):
        from yaml2pptx.components.panels import render_comparison
        render_comparison(
            blank_slide, theme, title="Compare",
            left_panel={"header": "BEFORE", "title": "Old", "rows": [{"label": "K", "value": "V"}]},
            right_panel={"header": "AFTER", "title": "New", "rows": [{"label": "K", "value": "V"}]},
            footer_text="Summary",
        )
        assert len(blank_slide.shapes) > 0


class TestSectionDivider:
    def test_renders(self, blank_slide, theme):
        from yaml2pptx.components.panels import render_section_divider
        render_section_divider(blank_slide, theme, number="03", title="Section", subtitle="Sub")
        assert len(blank_slide.shapes) > 0


class TestTimeline:
    def test_renders_phases(self, blank_slide, theme):
        from yaml2pptx.components.timeline import render_timeline
        phases = [
            {"label": "Q1", "title": "Start", "active": True, "items": ["A"]},
            {"label": "Q2", "title": "Middle", "description": "Desc"},
            {"label": "Q3", "title": "End"},
        ]
        render_timeline(blank_slide, theme, title="Timeline", phases=phases)
        assert len(blank_slide.shapes) > 0

    def test_empty_phases(self, blank_slide, theme):
        from yaml2pptx.components.timeline import render_timeline
        render_timeline(blank_slide, theme, title="Empty")


class TestProcess:
    def test_renders_steps(self, blank_slide, theme):
        from yaml2pptx.components.timeline import render_process
        steps = [
            {"title": "Step 1", "icon": "search", "description": "D1"},
            {"title": "Step 2", "items": ["A", "B"]},
            {"number": "3", "title": "Step 3"},
        ]
        render_process(blank_slide, theme, title="Process", steps=steps)
        assert len(blank_slide.shapes) > 0

    def test_empty_steps(self, blank_slide, theme):
        from yaml2pptx.components.timeline import render_process
        render_process(blank_slide, theme, title="Empty")


class TestQuote:
    def test_renders_light(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_quote
        render_quote(blank_slide, theme, quote="Test quote",
                     attribution="Speaker", source="Source")
        assert len(blank_slide.shapes) > 0

    def test_renders_dark(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_quote
        render_quote(blank_slide, theme, quote="Dark quote", dark=True)
        bg = blank_slide.background.fill
        assert bg.fore_color.rgb == theme.colors.dark_navy


class TestKeyMetrics:
    def test_renders_metrics(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_key_metrics
        metrics = [
            {"value": "99%", "label": "Uptime", "trend": "up", "trend_label": "+1%"},
            {"value": "42", "label": "Count", "trend": "down", "trend_label": "-5"},
            {"value": "7", "label": "Score", "color": "accent"},
        ]
        render_key_metrics(blank_slide, theme, title="Metrics", metrics=metrics)
        assert len(blank_slide.shapes) > 0

    def test_many_metrics_grid(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_key_metrics
        metrics = [{"value": str(i), "label": f"M{i}"} for i in range(6)]
        render_key_metrics(blank_slide, theme, metrics=metrics)

    def test_empty_metrics(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_key_metrics
        render_key_metrics(blank_slide, theme, title="Empty")


class TestTable:
    def test_renders_table(self, blank_slide, theme):
        from yaml2pptx.components.renderer import render_table
        render_table(blank_slide, theme, title="Data",
                     headers=["Name", "Value"],
                     rows=[["A", "1"], ["B", "2"]])
        assert len(blank_slide.shapes) > 0

    def test_renders_without_headers(self, blank_slide, theme):
        from yaml2pptx.components.renderer import render_table
        render_table(blank_slide, theme, title="Data",
                     rows=[["A", "1"], ["B", "2"]])

    def test_empty_table(self, blank_slide, theme):
        from yaml2pptx.components.renderer import render_table
        render_table(blank_slide, theme, title="Empty")

    def test_content_slide_with_table(self, tmp_path):
        slides = [{"type": "content", "title": "T",
                   "table": {"headers": ["H1", "H2"], "rows": [["a", "b"]]}}]
        out = render_presentation(slides, output_path=tmp_path / "out.pptx")
        assert out.exists()


class TestImageText:
    def test_renders_without_image(self, blank_slide, theme):
        from yaml2pptx.components.renderer import render_image_text
        render_image_text(blank_slide, theme, title="T", content=["Point 1"])
        assert len(blank_slide.shapes) > 0

    def test_renders_with_missing_image_placeholder(self, blank_slide, theme):
        from yaml2pptx.components.renderer import render_image_text
        render_image_text(blank_slide, theme, title="T",
                          image="nonexistent.png", content=["Point"])
        # Should render placeholder, not crash

    def test_image_position_right(self, blank_slide, theme):
        from yaml2pptx.components.renderer import render_image_text
        render_image_text(blank_slide, theme, title="T",
                          image_position="right", content=["Point"])


class TestChecklist:
    def test_renders_items(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_checklist
        items = [
            {"text": "Done task", "status": "done"},
            {"text": "In progress", "status": "in_progress", "note": "WIP"},
            {"text": "Pending", "status": "pending"},
            {"text": "Blocked", "status": "blocked", "note": "Waiting"},
        ]
        render_checklist(blank_slide, theme, title="Checklist", items=items)
        assert len(blank_slide.shapes) > 0

    def test_multi_column(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_checklist
        items = [{"text": f"Item {i}", "status": "pending"} for i in range(8)]
        render_checklist(blank_slide, theme, items=items, columns=2)

    def test_empty_items(self, blank_slide, theme):
        from yaml2pptx.components.quote import render_checklist
        render_checklist(blank_slide, theme, title="Empty")


# --- Icons ---

class TestIcons:
    def test_get_icon_known(self):
        from yaml2pptx.components.icons import get_icon
        char, font = get_icon("shield")
        assert char
        assert font

    def test_get_icon_unknown_fallback(self):
        from yaml2pptx.components.icons import get_icon
        char, font = get_icon("nonexistent")
        assert char == "●"
        assert font == "Calibri"

    def test_get_icon_case_insensitive(self):
        from yaml2pptx.components.icons import get_icon
        char1, _ = get_icon("Shield")
        char2, _ = get_icon("shield")
        assert char1 == char2

    def test_aliases(self):
        from yaml2pptx.components.icons import get_icon
        gpu_char, _ = get_icon("gpu")
        chip_char, _ = get_icon("chip")
        assert gpu_char == chip_char

    def test_add_icon_shape(self, blank_slide, theme):
        from yaml2pptx.components.icons import add_icon_shape
        add_icon_shape(blank_slide, 1.0, 1.0, 0.7, "shield", bg_color=theme.colors.primary)
        assert len(blank_slide.shapes) > 0

    def test_get_icon_names(self):
        from yaml2pptx.components.icons import get_icon_names
        names = get_icon_names()
        assert "shield" in names
        assert "brain" in names
        assert len(names) > 30


# --- resolve_color ---

class TestResolveColor:
    def test_none_returns_default(self):
        from pptx.dml.color import RGBColor
        default = RGBColor(0xFF, 0x00, 0x00)
        assert resolve_color(None, default) == default

    def test_empty_returns_default(self):
        from pptx.dml.color import RGBColor
        default = RGBColor(0xFF, 0x00, 0x00)
        assert resolve_color("", default) == default

    def test_hex_color(self):
        from pptx.dml.color import RGBColor
        default = RGBColor(0x00, 0x00, 0x00)
        result = resolve_color("#00A9A5", default)
        assert result == RGBColor(0x00, 0xA9, 0xA5)

    def test_named_color(self):
        from pptx.dml.color import RGBColor
        default = RGBColor(0x00, 0x00, 0x00)
        result = resolve_color("success", default)
        assert result == RGBColor(0x05, 0x96, 0x69)

    def test_unknown_name_returns_default(self):
        from pptx.dml.color import RGBColor
        default = RGBColor(0xFF, 0x00, 0x00)
        assert resolve_color("rainbow", default) == default


# --- Full YAML integration ---

class TestFullYamlBuild:
    def test_gpu_strategi(self, tmp_path):
        yaml_path = Path("examples/gpu_strategi.yaml")
        if not yaml_path.exists():
            pytest.skip("Example file not found")
        text = yaml_path.read_text(encoding="utf-8")
        data = yaml.safe_load(text)
        out = render_presentation(
            data["slides"],
            output_path=tmp_path / "gpu.pptx",
            metadata=data.get("metadata"),
        )
        prs = Presentation(str(out))
        assert len(prs.slides) == 10

    def test_all_slide_types_in_one(self, tmp_path):
        """Build a presentation with every slide type to catch regressions."""
        slides = [
            {"type": "title_page", "title": "Test", "subtitle": "Sub"},
            {"type": "agenda", "title": "Agenda", "items": [{"title": "T"}]},
            {"type": "stat_cards", "title": "Stats", "cards": [{"stat": "1"}]},
            {"type": "definition_cards", "title": "Defs", "cards": [{"term": "X"}]},
            {"type": "content_cards", "title": "Cards", "cards": [{"title": "C"}]},
            {"type": "icon_cards", "message": "Msg", "cards": [{"title": "I"}]},
            {"type": "two_panels", "title": "Panels",
             "left_panel": {"title": "L"}, "right_panel": {"title": "R"}},
            {"type": "comparison", "title": "Comp",
             "left_panel": {"header": "A", "rows": []},
             "right_panel": {"header": "B", "rows": []}},
            {"type": "section_divider", "number": "01", "title": "Div"},
            {"type": "timeline", "title": "TL", "phases": [{"title": "P1"}]},
            {"type": "process", "title": "Proc", "steps": [{"title": "S1"}]},
            {"type": "quote", "quote": "Q"},
            {"type": "key_metrics", "title": "KM", "metrics": [{"value": "1", "label": "L"}]},
            {"type": "checklist", "title": "CL", "items": [{"text": "T", "status": "done"}]},
            {"type": "table", "title": "Tbl", "headers": ["A"], "rows": [["1"]]},
            {"type": "image_text", "title": "ImgTxt", "content": ["Point"]},
            {"type": "content", "title": "Fallback", "content": ["Line"]},
        ]
        out = render_presentation(slides, output_path=tmp_path / "all.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == len(slides)
